import os
import re
import json
import requests
import threading
import contextlib
import urllib3.util.connection as conn
from urllib.parse import urlparse, urljoin
from bs4 import BeautifulSoup
from duckduckgo_search import DDGS
from vcdiligence.logging_config import logger

_thread_local = threading.local()
playwright_semaphore = threading.Semaphore(2)

def _get_pinned_ips():
    if not hasattr(_thread_local, "pinned_ips"):
        _thread_local.pinned_ips = {}
    return _thread_local.pinned_ips

orig_create_connection = conn.create_connection

def custom_create_connection(address, *args, **kwargs):
    host, port = address
    pinned_ips = _get_pinned_ips()
    host_lower = host.lower() if host else ""
    if host_lower in pinned_ips:
        return orig_create_connection((pinned_ips[host_lower], port), *args, **kwargs)
    if host_lower.startswith("www."):
        parent_host = host_lower[4:]
        if parent_host in pinned_ips:
            return orig_create_connection((pinned_ips[parent_host], port), *args, **kwargs)
    return orig_create_connection(address, *args, **kwargs)

conn.create_connection = custom_create_connection

@contextlib.contextmanager
def pinned_dns(domain, ip):
    if not ip:
        yield
        return
    pinned_ips = _get_pinned_ips()
    domain_lower = domain.lower()
    pinned_ips[domain_lower] = ip
    try:
        yield
    finally:
        pinned_ips.pop(domain_lower, None)

class SmartScraper:
    @staticmethod
    def get_domain(url):
        parsed = urlparse(url)
        domain = parsed.netloc or parsed.path
        if domain.startswith("www."):
            domain = domain[4:]
        return domain

    @staticmethod
    def clean_text(text):
        text = re.sub(r'\s+', ' ', text)
        return text.strip()

    @classmethod
    def scrape_with_requests(cls, url, validated_ip=None):
        domain = cls.get_domain(url)
        with pinned_dns(domain, validated_ip):
            try:
                headers = {
                    "User-Agent": "Mozilla/5.0 (iPhone; CPU iPhone OS 15_0 like Mac OS X) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/15.0 Mobile/15E148 Safari/604.1"
                }
                response = requests.get(url, headers=headers, timeout=12)
                if response.status_code != 200:
                    logger.warning(f"Requests scrape failed with HTTP {response.status_code} for {url}")
                    return None

                soup = BeautifulSoup(response.text, "html.parser")
                for script_or_style in soup(["script", "style", "nav", "footer"]):
                    script_or_style.decompose()

                cleaned = cls.clean_text(soup.get_text())
                if len(cleaned) < 300:
                    logger.warning(f"Requests scrape returned very short content ({len(cleaned)} chars) for {url}")
                    return None
                return cleaned
            except Exception as e:
                logger.error(f"Error scraping with requests on {url}: {str(e)}")
                return None

    @classmethod
    def scrape_with_playwright(cls, url, validated_ip=None):
        logger.info(f"Using Playwright headless fallback for {url}")
        try:
            from playwright.sync_api import sync_playwright
            with sync_playwright() as p:
                domain = cls.get_domain(url)
                args = []
                if validated_ip:
                    args.append(f"--host-resolver-rules=MAP {domain} {validated_ip}, MAP www.{domain} {validated_ip}")
                with playwright_semaphore:
                    browser = p.chromium.launch(headless=True, args=args)
                    # Create a context with custom User-Agent
                    context = browser.new_context(
                        user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
                    )
                    page = context.new_page()
                    page.goto(url, timeout=20000, wait_until="load")
                    # Wait 1s for any dynamically-rendered text
                    page.wait_for_timeout(1000)
                    content = page.content()
                    browser.close()

                soup = BeautifulSoup(content, "html.parser")
                for script_or_style in soup(["script", "style", "nav", "footer"]):
                    script_or_style.decompose()

                cleaned = cls.clean_text(soup.get_text())
                if len(cleaned) < 100:
                    return f"[Could not verify content for {url} - page load returned insufficient text]"
                return cleaned
        except Exception as e:
            logger.error(f"Playwright fallback also failed for {url}: {str(e)}")
            return f"[Could not verify content for {url} due to connection error or security block]"

    @classmethod
    def scrape_url(cls, url, validated_ip=None):
        # First try requests
        text = cls.scrape_with_requests(url, validated_ip=validated_ip)
        if text:
            return text
        # If requests fails or returns very short text, try Playwright
        return cls.scrape_with_playwright(url, validated_ip=validated_ip)

    @classmethod
    def get_internal_links(cls, base_url, validated_ip=None):
        links = set()
        domain = cls.get_domain(base_url)
        with pinned_dns(domain, validated_ip):
            try:
                headers = {
                    "User-Agent": "Mozilla/5.0 (iPhone; CPU iPhone OS 15_0 like Mac OS X) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/15.0 Mobile/15E148 Safari/604.1"
                }
                # Attempt to fetch links using requests, fallback to Playwright if needed
                response = requests.get(base_url, headers=headers, timeout=10)
                html = response.text if response.status_code == 200 else ""
                if not html or len(html) < 2000:
                    logger.info(f"Using Playwright headless fallback to find internal links for {base_url}")
                    try:
                        from playwright.sync_api import sync_playwright
                        with sync_playwright() as p:
                            args = []
                            if validated_ip:
                                args.append(f"--host-resolver-rules=MAP {domain} {validated_ip}, MAP www.{domain} {validated_ip}")
                            with playwright_semaphore:
                                browser = p.chromium.launch(headless=True, args=args)
                                page = browser.new_page()
                                page.goto(base_url, timeout=15000, wait_until="load")
                                html = page.content()
                                browser.close()
                    except Exception:
                        pass

                if html:
                    soup = BeautifulSoup(html, "html.parser")
                    domain = cls.get_domain(base_url)
                    for anchor in soup.find_all("a", href=True):
                        href = anchor["href"]
                        full_url = urljoin(base_url, href)
                        parsed_full = urlparse(full_url)
                        full_domain = parsed_full.netloc
                        if full_domain.startswith("www."):
                            full_domain = full_domain[4:]

                        if full_domain == domain:
                            path_lower = parsed_full.path.lower()
                            if any(kw in path_lower for kw in ["about", "team", "pricing", "product", "features", "career", "contact"]):
                                links.add(full_url)
            except Exception as e:
                logger.error(f"Error getting internal links for {base_url}: {str(e)}")
            return list(links)[:4]

    @classmethod
    def search_duckduckgo(cls, query, count=3):
        results = []
        try:
            with DDGS() as ddgs:
                for result in ddgs.text(query, max_results=count):
                    results.append({
                        "title": result.get("title", ""),
                        "link": result.get("href", ""),
                        "snippet": result.get("body", "")
                    })
        except Exception as e:
            logger.error(f"DuckDuckGo search error for query '{query}': {str(e)}")

        # If absolutely no results, return explicit trace to prevent LLM assuming details
        if not results:
            return [{"title": "No public search results", "link": "", "snippet": f"[No public search records found under specific query: {query}]"}]
        return results

    @classmethod
    def analyze_startup(cls, url, validated_ip=None):
        if not validated_ip:
            from vcdiligence.validator import validate_url_for_ssrf
            try:
                url, validated_ip = validate_url_for_ssrf(url)
            except Exception as e:
                logger.error(f"SSRF validation failed for {url} inside analyze_startup: {str(e)}")
                raise

        domain = cls.get_domain(url)
        cache_dir = os.path.join(os.path.dirname(__file__), "cache")
        os.makedirs(cache_dir, exist_ok=True)
        cache_path = os.path.join(cache_dir, f"{domain}.json")

        # Check local cache first
        if os.path.exists(cache_path):
            try:
                with open(cache_path, "r", encoding="utf-8") as file:
                    return json.load(file)
            except Exception:
                pass

        logger.info(f"Starting analysis for startup: {url}")
        homepage_content = cls.scrape_url(url, validated_ip=validated_ip)
        company_name = domain.split('.')[0].capitalize()

        internal_content = {}
        internal_links = cls.get_internal_links(url, validated_ip=validated_ip)
        for link in internal_links:
            link_path = urlparse(link).path
            internal_content[link_path] = cls.scrape_url(link, validated_ip=validated_ip)[:1500]

        # Explicitly record missing sub-pages so Omission Analyst is aware
        expected_keywords = ["pricing", "team", "about", "features"]
        found_keywords = [kw for kw in expected_keywords if any(kw in path.lower() for path in internal_content.keys())]
        for kw in expected_keywords:
            if kw not in found_keywords:
                internal_content[f"/{kw}-missing-page"] = f"[Could not verify {kw} details: no dedicated /{kw} page found or loaded]"

        ddg_results = {}
        search_queries = {
            "competitors": f"{company_name} competitors alternative SaaS",
            "team_and_founders": f"{company_name} founders team LinkedIn",
            "market_and_funding": f"{company_name} Crunchbase funding traction",
            "pricing_and_product": f"{company_name} pricing product reviews"
        }

        for category, query in search_queries.items():
            ddg_results[category] = cls.search_duckduckgo(query)

        analysis_payload = {
            "company_name": company_name,
            "company_url": url,
            "homepage_summary": homepage_content[:3000],
            "internal_pages": internal_content,
            "search_insights": ddg_results
        }

        try:
            with open(cache_path, "w", encoding="utf-8") as file:
                json.dump(analysis_payload, file, ensure_ascii=False, indent=2)
        except Exception:
            pass

        return analysis_payload

    @classmethod
    def scrape_linkedin(cls, linkedin_url):
        logger.info(f"Extracting LinkedIn context from URL: {linkedin_url}")
        results = cls.search_duckduckgo(linkedin_url, count=2)
        snippet_text = ""
        for r in results:
            snippet_text += f"\nTitle: {r.get('title', '')}\nSnippet: {r.get('snippet', '')}\n"
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
            }
            res = requests.get(linkedin_url, headers=headers, timeout=8)
            if res.status_code == 200:
                soup = BeautifulSoup(res.text, "html.parser")
                og_title = soup.find("meta", property="og:title")
                og_desc = soup.find("meta", property="og:description")
                meta_desc = soup.find("meta", {"name": "description"})
                title_val = og_title["content"] if og_title else ""
                desc_val = og_desc["content"] if og_desc else (meta_desc["content"] if meta_desc else "")
                if title_val or desc_val:
                    snippet_text += f"\nLinkedIn OG Title: {title_val}\nLinkedIn OG Description: {desc_val}\n"
        except Exception as e:
            logger.warning(f"Direct LinkedIn requests scraping failed or blocked: {str(e)}")

        inferred_url = None
        urls = re.findall(r'https?://(?:www\.)?([a-zA-Z0-9-]+(?:\.[a-zA-Z0-9-]+)+)', snippet_text)
        for u in urls:
            if not any(ignored in u.lower() for ignored in ["linkedin.com", "google.com", "twitter.com", "facebook.com", "crunchbase.com", "wikipedia.org"]):
                inferred_url = f"https://{u}"
                break

        company_name = None
        match = re.search(r'/company/([a-zA-Z0-9-]+)', linkedin_url)
        if match:
            company_name = match.group(1).replace("-", " ")
        else:
            match_in = re.search(r'/in/([a-zA-Z0-9-]+)', linkedin_url)
            if match_in:
                company_name = match_in.group(1).replace("-", " ")

        if company_name and not inferred_url:
            search_res = cls.search_duckduckgo(f"{company_name} official website", count=1)
            if search_res and search_res[0].get("link"):
                inferred_url = search_res[0].get("link")

        return {
            "linkedin_data": snippet_text.strip() or f"[No public LinkedIn snippets found for URL: {linkedin_url}]",
            "inferred_url": inferred_url,
            "company_name": company_name
        }

    @classmethod
    def extract_text_from_pdf(cls, file_path):
        try:
            import pypdf
            reader = pypdf.PdfReader(file_path)
            text = ""
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    text += t + "\n"
            return text
        except Exception as e:
            logger.error(f"Error reading PDF: {str(e)}")
            return ""

    @classmethod
    def extract_text_from_pptx(cls, file_path):
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error reading PPTX: {str(e)}")
            return ""

    @classmethod
    def extract_url_from_text(cls, text):
        urls = re.findall(r'(https?://[^\s()<>]+(?:\.[^\s()<>]+)+)', text)
        if not urls:
            urls = re.findall(r'\b(?:[a-zA-Z0-9-]+\.)+(?:com|io|co|net|org|ai|edu|gov|mx|lat|app)\b', text)
        for u in urls:
            u_lower = u.lower()
            if not any(generic in u_lower for generic in ["schema.org", "w3.org", "adobe.com", "microsoft.com", "pypdf", "openxmlformats"]):
                if not u_lower.startswith("http"):
                    u = f"https://{u}"
                return u
        return None
